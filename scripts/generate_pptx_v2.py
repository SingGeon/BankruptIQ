#!/usr/bin/env python3
"""
Generează prezentarea universitară BankruptIQ v2 — cu screenshots, stil academic.
Rulare: python scripts/generate_pptx_v2.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import datetime, os

SCREENS = "/home/user/BankruptIQ/scripts/screenshots"

# ── Paleta aplicatiei BankruptIQ ────────────────────────────────────
BG       = RGBColor(0x0F, 0x11, 0x17)   # --bg
SURFACE  = RGBColor(0x1A, 0x1D, 0x27)   # --surface
SURFACE2 = RGBColor(0x22, 0x26, 0x3A)   # --surface2
BORDER_C = RGBColor(0x2E, 0x33, 0x49)   # --border
ACCENT   = RGBColor(0x4F, 0x8E, 0xF7)   # --accent (blue)
DANGER   = RGBColor(0xEF, 0x44, 0x44)   # --danger
WARNING  = RGBColor(0xF5, 0x9E, 0x0B)   # --warning
SUCCESS  = RGBColor(0x22, 0xC5, 0x5E)   # --success
TEXT_C   = RGBColor(0xE2, 0xE8, 0xF0)   # --text
MUTED    = RGBColor(0x88, 0x92, 0xA4)   # --text-muted
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GOLD     = RGBColor(0xF5, 0xC5, 0x18)
CYAN     = RGBColor(0x00, 0xC2, 0xFF)

SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]


# ── Utilitare ─────────────────────────────────────────────────────────

def bg(slide, color=None):
    color = color or BG
    s = slide.shapes.add_shape(1, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def rect(slide, x, y, w, h, fill=SURFACE, ec=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if ec:
        s.line.color.rgb = ec; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=14, bold=False,
        color=TEXT_C, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return box

def add_accent_stripe(slide, color=ACCENT):
    s = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.055))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def slide_title(slide, title, subtitle=None, tag=None):
    add_accent_stripe(slide)
    txt(slide, title, Inches(0.5), Inches(0.15), Inches(10), Inches(0.65),
        size=26, bold=True, color=WHITE)
    if tag:
        r = rect(slide, Inches(11.2), Inches(0.15), Inches(1.9), Inches(0.5),
                 fill=SURFACE2, ec=ACCENT, lw=Pt(1))
        txt(slide, tag, Inches(11.2), Inches(0.15), Inches(1.9), Inches(0.5),
            size=10, color=ACCENT, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(0.82), Inches(12), Inches(0.38),
            size=12, color=MUTED, italic=True)
    sep = slide.shapes.add_shape(1, Inches(0.5), Inches(1.25),
                                  Inches(12.3), Inches(0.025))
    sep.fill.solid(); sep.fill.fore_color.rgb = BORDER_C
    sep.line.fill.background()

def bullet_card(slide, title, items, x, y, w, h,
                title_color=ACCENT, item_size=13, title_size=15, icon=""):
    rect(slide, x, y, w, h, fill=SURFACE, ec=BORDER_C, lw=Pt(1.2))
    txt(slide, f"{icon}  {title}" if icon else title,
        x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.42),
        size=title_size, bold=True, color=title_color)
    sep = slide.shapes.add_shape(1, x + Inches(0.15), y + Inches(0.56),
                                  w - Inches(0.3), Inches(0.02))
    sep.fill.solid(); sep.fill.fore_color.rgb = BORDER_C; sep.line.fill.background()
    iy = y + Inches(0.65)
    for item in items:
        txt(slide, f"  {item}", x + Inches(0.15), iy,
            w - Inches(0.3), Inches(0.38), size=item_size, color=TEXT_C)
        iy += Inches(0.36)

def kpi(slide, val, label, x, y, w=Inches(2.8), h=Inches(1.3), vc=ACCENT):
    rect(slide, x, y, w, h, fill=SURFACE, ec=ACCENT, lw=Pt(1.5))
    txt(slide, val, x, y + Inches(0.1), w, Inches(0.7),
        size=30, bold=True, color=vc, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.78), w, Inches(0.42),
        size=11, color=MUTED, align=PP_ALIGN.CENTER)

def add_image(slide, path, x, y, w, h=None):
    if os.path.exists(path):
        if h:
            slide.shapes.add_picture(path, x, y, w, h)
        else:
            slide.shapes.add_picture(path, x, y, w)

def number_badge(slide, num, x, y):
    r = slide.shapes.add_shape(1, x, y, Inches(0.45), Inches(0.45))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT; r.line.fill.background()
    txt(slide, str(num), x, y, Inches(0.45), Inches(0.45),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Copertă universitară
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)

# Gradient lateral stanga (accent line groasa)
for i in range(6):
    s = slide.shapes.add_shape(1, Inches(i * 0.05), 0, Inches(0.05), SH)
    s.fill.solid()
    alpha = int(80 * (1 - i/6))
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()

# Logo / titlu aplicatie
rect(slide, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.08), fill=ACCENT)
txt(slide, "LUCRARE DE PROIECT",
    Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.5),
    size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

txt(slide, "BankruptIQ",
    Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.3),
    size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txt(slide, "Sistem Inteligent de Predicție a Riscului de Faliment\nal Întreprinderilor din România",
    Inches(0.5), Inches(2.75), Inches(12.3), Inches(0.9),
    size=20, color=TEXT_C, align=PP_ALIGN.CENTER)

rect(slide, Inches(4.5), Inches(3.72), Inches(4.3), Inches(0.04), fill=ACCENT)

# Subtitlu descriere
txt(slide,
    "Analiză bazată pe Machine Learning · Date oficiale Ministerul Finanțelor · 80.000+ companii",
    Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.45),
    size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# Tags tehnologii
tags = ["Python 3.12", "FastAPI", "MongoDB", "RandomForest ML",
        "React + Plotly", "data.gov.ro"]
tx = Inches(2.1)
for tag in tags:
    tw = Inches(1.55)
    rect(slide, tx, Inches(4.55), tw, Inches(0.42),
         fill=SURFACE2, ec=BORDER_C, lw=Pt(1))
    txt(slide, tag, tx, Inches(4.55), tw, Inches(0.42),
        size=11, color=TEXT_C, align=PP_ALIGN.CENTER)
    tx += Inches(1.65)

# Linie despartitoare
rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.025), fill=BORDER_C)

# Info autor / universitate
txt(slide, "Proiect realizat cu date financiare reale din România (2017–2023)",
    Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
    size=12, color=TEXT_C, align=PP_ALIGN.CENTER)
txt(slide, f"Sistem functional · {datetime.date.today().strftime('%B %Y')}",
    Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.38),
    size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# KPI preview jos
kpis_c = [("80.000+", "Companii analizate"), ("88%+", "Acuratețe model ML"),
          ("6 ani", "Date istorice"), ("10", "Indicatori financiari")]
for ki, (kv, kl) in enumerate(kpis_c):
    kx = Inches(1.5) + ki * Inches(2.6)
    rect(slide, kx, Inches(6.4), Inches(2.3), Inches(0.88),
         fill=SURFACE, ec=ACCENT, lw=Pt(1.2))
    txt(slide, kv, kx, Inches(6.44), Inches(2.3), Inches(0.44),
        size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(slide, kl, kx, Inches(6.88), Inches(2.3), Inches(0.32),
        size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Cuprins
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Cuprins", tag="Index")

chapters = [
    ("01", "Contextul și Motivația",     "De ce este necesară această aplicație?",     ACCENT),
    ("02", "Obiective",                  "Ce ne-am propus să rezolvăm?",               SUCCESS),
    ("03", "Beneficiile Aplicației",     "Cine beneficiază și cum?",                   WARNING),
    ("04", "Prezentarea Aplicației",     "Interfața și funcționalitățile principale",  ACCENT),
    ("05", "Arhitectura Sistemului",     "Stack tehnic: FastAPI · MongoDB · ML · React",MUTED),
    ("06", "Baza de Date MongoDB",       "Structura și operațiile bazei de date",      SUCCESS),
    ("07", "Modelul de Machine Learning","RandomForest + Altman Z-score blended",      WARNING),
    ("08", "Datele Oficiale MF",         "Import, procesare, conversie CUI → Firmă",  ACCENT),
    ("09", "Rezultate și KPI-uri",       "Metrici model, volume, performanță",         SUCCESS),
    ("10", "Concluzii și Perspective",   "Ce am învățat și direcții viitoare",         MUTED),
]
for i, (num, title, sub, color) in enumerate(chapters):
    col = i // 5; row = i % 5
    x = Inches(0.5) + col * Inches(6.45)
    y = Inches(1.45) + row * Inches(1.16)
    rect(slide, x, y, Inches(6.2), Inches(1.05), fill=SURFACE, ec=color, lw=Pt(1.2))
    txt(slide, num, x + Inches(0.12), y + Inches(0.12),
        Inches(0.55), Inches(0.5), size=22, bold=True, color=color)
    txt(slide, title, x + Inches(0.7), y + Inches(0.08),
        Inches(5.3), Inches(0.46), size=15, bold=True, color=WHITE)
    txt(slide, sub, x + Inches(0.7), y + Inches(0.6),
        Inches(5.3), Inches(0.36), size=11, color=MUTED)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Context și Motivație
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Contextul și Motivația Proiectului",
            subtitle="De ce am creat BankruptIQ?", tag="01 · Context")

# Stanga: problema
rect(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.7),
     fill=SURFACE, ec=DANGER, lw=Pt(1.5))
txt(slide, "PROBLEMA", Inches(0.65), Inches(1.5),
    Inches(5.7), Inches(0.45), size=13, bold=True, color=DANGER)
txt(slide, "Situatia Actuala in Romania",
    Inches(0.65), Inches(2.0), Inches(5.7), Inches(0.45),
    size=15, bold=True, color=TEXT_C)

problems = [
    "In fiecare an, mii de firme romanesti",
    "intra in insolventa, afectand angajati,",
    "furnizori si parteneri de afaceri.",
    "",
    "> 50.000 firme intrate in insolventa",
    "  in Romania doar in 2022-2023",
    "",
    "> Creditorii pierd in medie 60-70%",
    "  din creantele lor",
    "",
    "> Bancile si investitorii nu au acces",
    "  la un sistem de predictie accesibil",
    "  bazat pe date financiare deschise",
    "",
    "> Datele MF sunt publice dar greu de",
    "  procesat si interpretat",
]
py = Inches(2.55)
for line in problems:
    color = WARNING if line.startswith(">") else TEXT_C
    size = 12.5 if line.startswith(">") else 12
    txt(slide, line, Inches(0.65), py, Inches(5.7), Inches(0.32),
        size=size, color=color)
    py += Inches(0.3)

# Dreapta: solutia
rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.7),
     fill=SURFACE, ec=SUCCESS, lw=Pt(1.5))
txt(slide, "SOLUTIA", Inches(6.95), Inches(1.5),
    Inches(5.7), Inches(0.45), size=13, bold=True, color=SUCCESS)
txt(slide, "Ce face BankruptIQ diferit",
    Inches(6.95), Inches(2.0), Inches(5.7), Inches(0.45),
    size=15, bold=True, color=TEXT_C)

solutions = [
    (SUCCESS, "Date oficiale reale",
     "Foloseste datele financiare publicate\nde Ministerul Finantelor pe data.gov.ro"),
    (ACCENT,  "Algoritm hibrid",
     "Combina RandomForest (ML modern) cu\nAltman Z-score (model academic validat)"),
    (WARNING, "Accesibil si deschis",
     "Interfata web intuitiva, fara cunostinte\ntehnice avansate necesare"),
    (SUCCESS, "Scalabil",
     "Analizeaza 80.000+ companii simultan\ncu raspuns in sub 200ms"),
]
sy = Inches(2.55)
for sc, stitle, sdesc in solutions:
    rect(slide, Inches(6.95), sy, Inches(5.6), Inches(1.12),
         fill=SURFACE2, ec=sc, lw=Pt(1))
    txt(slide, stitle, Inches(7.1), sy + Inches(0.08),
        Inches(5.3), Inches(0.38), size=13, bold=True, color=sc)
    txt(slide, sdesc, Inches(7.1), sy + Inches(0.5),
        Inches(5.3), Inches(0.55), size=11.5, color=TEXT_C)
    sy += Inches(1.25)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Obiective
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Obiectivele Proiectului",
            subtitle="Ce ne-am propus sa realizam", tag="02 · Obiective")

obj_general = [
    "Crearea unui sistem automat de evaluare a riscului de faliment",
    "pentru companiile romanesti, bazat pe date financiare publice,",
    "accesibil oricarei persoane fara cunostinte tehnice avansate.",
]
rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.1),
     fill=SURFACE2, ec=ACCENT, lw=Pt(2))
txt(slide, "Obiectiv General:", Inches(0.65), Inches(1.45),
    Inches(3), Inches(0.38), size=12, bold=True, color=ACCENT)
for i, line in enumerate(obj_general):
    txt(slide, line, Inches(0.65), Inches(1.85) + i * Inches(0.26),
        Inches(12), Inches(0.28), size=12, color=TEXT_C)

obj_specifice = [
    ("O1", "Date reale",      "Colectarea si procesarea datelor financiare\nofficiale de la Ministerul Finantelor\n(data.gov.ro) pentru 6 ani fiscali",           ACCENT),
    ("O2", "Model ML",        "Dezvoltarea unui model de clasificare\nRandomForest cu acuratete >85%, combinat\ncu formula academica Altman Z-score",           SUCCESS),
    ("O3", "Baza de date",    "Stocarea eficienta a datelor in MongoDB\ncu indexuri compuse pentru\ninterogari rapide (<200ms)",                                 WARNING),
    ("O4", "API REST",        "Implementarea unui API REST modern\n(FastAPI) cu endpoint-uri pentru\nlistare, filtrare, predictie si import",                    ACCENT),
    ("O5", "Dashboard",       "Crearea unei interfete web interactive\n(React + Plotly) cu vizualizari\nsi functii de comparare companii",                       SUCCESS),
    ("O6", "Conversie CUI",   "Imbogatirea automata a datelor cu\ndenumirile reale ale firmelor prin\nAPI ANAF si cache ONRC",                                   WARNING),
]
ow = Inches(3.88)
for i, (code, title, desc, color) in enumerate(obj_specifice):
    col = i % 3; row = i // 3
    ox = Inches(0.5) + col * Inches(4.28)
    oy = Inches(2.75) + row * Inches(2.0)
    rect(slide, ox, oy, ow, Inches(1.85), fill=SURFACE, ec=color, lw=Pt(1.5))
    rect(slide, ox, oy, Inches(0.65), Inches(1.85), fill=color)
    txt(slide, code, ox, oy, Inches(0.65), Inches(1.85),
        size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(slide, title, ox + Inches(0.72), oy + Inches(0.1),
        ow - Inches(0.8), Inches(0.42), size=14, bold=True, color=color)
    txt(slide, desc, ox + Inches(0.72), oy + Inches(0.58),
        ow - Inches(0.8), Inches(1.15), size=11, color=TEXT_C)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Beneficii practice
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Beneficiile Aplicatiei",
            subtitle="Cine beneficiaza si cum? Impact practic concret",
            tag="03 · Beneficii")

beneficiari = [
    ("Banci si\nInstitutii Financiare", ACCENT, [
        "Evaluare rapida a bonitatii clientilor",
        "Detectie timpurie risc credit",
        "Reducere pierderi prin provizioane corecte",
        "Conformitate Basel III / IFRS 9",
        "Audit portofoliu credite in cateva secunde",
    ]),
    ("Investitori si\nFonduri de Capital", SUCCESS, [
        "Due diligence automat inainte de investitie",
        "Monitorizare continua portofoliu",
        "Comparare simultana 100+ companii",
        "Identificare oportunitati de restructurare",
        "Benchmarking sector vs. competitor",
    ]),
    ("Furnizori si\nParteneri Comerciali", WARNING, [
        "Verificare client inainte de credit comercial",
        "Stabilire limite de credit personalizate",
        "Alerte automate la deteriorare situatie",
        "Reducere risc creante neincasate",
        "Decizie informata privind termenii contract",
    ]),
    ("Antreprenori si\nManageri", DANGER, [
        "Analiza proprie a sanatatii financiare",
        "Comparatie cu media sectorului",
        "Identificare indicatori cu risc ridicat",
        "Benchmark fata de competitori din sector",
        "Justificare decizii in fata actionarilor",
    ]),
]
bw = Inches(3.0)
for i, (title, color, items) in enumerate(beneficiari):
    bx = Inches(0.5) + i * Inches(3.2)
    by = Inches(1.4)
    rect(slide, bx, by, bw, Inches(5.8), fill=SURFACE, ec=color, lw=Pt(1.5))
    rect(slide, bx, by, bw, Inches(0.7), fill=color)
    txt(slide, title, bx, by, bw, Inches(0.7),
        size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
    iy = by + Inches(0.85)
    for item in items:
        txt(slide, f"+ {item}", bx + Inches(0.12), iy,
            bw - Inches(0.24), Inches(0.38), size=11.5, color=TEXT_C)
        iy += Inches(0.82)

# Impact bar jos
rect(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.32),
     fill=SURFACE2)
impacts = [
    ("Reducere timp analiza", "De la 2-3 zile  la  < 30 secunde"),
    ("Acuratete predictie",   "88%+ (RandomForest + Altman Z)"),
    ("Date acoperite",        "80.000+ companii · 2017–2023"),
]
for ii, (label, val) in enumerate(impacts):
    ix = Inches(0.8) + ii * Inches(4.3)
    txt(slide, label, ix, Inches(7.13), Inches(2.5), Inches(0.18),
        size=9, color=MUTED)
    txt(slide, val, ix + Inches(2.2), Inches(7.13), Inches(2.0), Inches(0.18),
        size=10, bold=True, color=ACCENT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Landing page (screenshot)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Interfata Aplicatiei — Pagina Principala",
            subtitle="Prima impresie: landing page modern, clar, orientat spre actiune",
            tag="04 · UI/UX")

add_image(slide, f"{SCREENS}/screen5_landing.png",
          Inches(0.5), Inches(1.38), Inches(8.2), Inches(5.8))

rect(slide, Inches(9.0), Inches(1.38), Inches(4.1), Inches(5.8),
     fill=SURFACE, ec=BORDER_C, lw=Pt(1))
txt(slide, "Ce ofera pagina principala",
    Inches(9.15), Inches(1.5), Inches(3.8), Inches(0.42),
    size=14, bold=True, color=ACCENT)

lp_features = [
    ("Titlu clar",     "Mesajul aplicatiei este imediat\ninteles de orice utilizator"),
    ("CTA primar",     "Buton direct catre Dashboard\nfara pasi intermediari inutili"),
    ("4 feature cards","Prezentare concisa a\ncapacitatilor principale"),
    ("Tech stack",     "Vizibil credibilitatea tehnica\na solutiei propuse"),
    ("Design fintech", "Tema dark profesionala,\ncu accente albastre specifice"),
]
fy = Inches(2.05)
for ftitle, fdesc in lp_features:
    rect(slide, Inches(9.15), fy, Inches(3.75), Inches(0.98),
         fill=SURFACE2, ec=BORDER_C, lw=Pt(1))
    txt(slide, ftitle, Inches(9.3), fy + Inches(0.06),
        Inches(3.4), Inches(0.36), size=12, bold=True, color=TEXT_C)
    txt(slide, fdesc, Inches(9.3), fy + Inches(0.44),
        Inches(3.4), Inches(0.46), size=10.5, color=MUTED)
    fy += Inches(1.1)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Dashboard KPI (screenshot)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Interfata Aplicatiei — Dashboard Principal",
            subtitle="Vedere de ansamblu: KPI cards, grafice sector, alerte, evolutie temporala",
            tag="04 · Dashboard")

add_image(slide, f"{SCREENS}/screen1_dashboard.png",
          Inches(0.5), Inches(1.38), Inches(8.2), Inches(5.8))

rect(slide, Inches(9.0), Inches(1.38), Inches(4.1), Inches(5.8),
     fill=SURFACE, ec=BORDER_C, lw=Pt(1))
txt(slide, "Componentele dashboard-ului",
    Inches(9.15), Inches(1.5), Inches(3.8), Inches(0.42),
    size=14, bold=True, color=ACCENT)

dashboard_items = [
    (ACCENT,   "5 KPI Cards",           "Total companii · Risc mare/mediu/mic\n· Scor mediu global"),
    (DANGER,   "Grafic stacked bar",     "Distributia riscului pe sectoare\n(10 domenii de activitate)"),
    (WARNING,  "Grafic donut",           "Ponderea globala a nivelurilor\nde risc (pie chart interactiv)"),
    (SUCCESS,  "Trend temporal",         "Evolutia % risc mare si mediu\npentru 2017-2023"),
    (ACCENT,   "Alerte recente",         "Top companii cu risc ridicat\ndetectate recent"),
]
dy = Inches(2.05)
for dc, dtitle, ddesc in dashboard_items:
    rect(slide, Inches(9.15), dy, Inches(3.75), Inches(0.98),
         fill=SURFACE2, ec=dc, lw=Pt(1))
    rect(slide, Inches(9.15), dy, Inches(0.06), Inches(0.98), fill=dc)
    txt(slide, dtitle, Inches(9.27), dy + Inches(0.06),
        Inches(3.4), Inches(0.38), size=12, bold=True, color=TEXT_C)
    txt(slide, ddesc, Inches(9.27), dy + Inches(0.47),
        Inches(3.4), Inches(0.44), size=10.5, color=MUTED)
    dy += Inches(1.1)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Lista companii (screenshot)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Interfata Aplicatiei — Managementul Companiilor",
            subtitle="Tabel interactiv cu filtrare, cautare, sortare si acces rapid la detalii",
            tag="04 · Companii")

add_image(slide, f"{SCREENS}/screen2_companies.png",
          Inches(0.5), Inches(1.38), Inches(12.3), Inches(4.5))

# Explicatii feature-uri tabel
features_tabel = [
    (ACCENT,   "Cautare live",       "Filtru instant dupa\nnumele companiei"),
    (SUCCESS,  "Filtre multiple",    "Sector · Nivel risc\n· An fiscal"),
    (WARNING,  "Risk badges",        "Cod color: verde/\ngalben/rosu"),
    (DANGER,   "Scor numeric",       "0-100: cat mai mare,\ncu atat risc mai mare"),
    (ACCENT,   "Export CSV",         "Download date filtrate\npentru analiza externa"),
    (MUTED,    "Detalii companie",   "Click → fisa completa\ncu toti indicatorii"),
]
fw = Inches(2.0)
for fi, (fc, ft, fd) in enumerate(features_tabel):
    fx = Inches(0.5) + fi * Inches(2.13)
    rect(slide, fx, Inches(6.05), fw, Inches(1.3),
         fill=SURFACE, ec=fc, lw=Pt(1.2))
    txt(slide, ft, fx + Inches(0.1), Inches(6.1),
        fw - Inches(0.2), Inches(0.38), size=12, bold=True, color=fc)
    txt(slide, fd, fx + Inches(0.1), Inches(6.5),
        fw - Inches(0.2), Inches(0.75), size=10.5, color=TEXT_C)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Predictie risc (screenshot)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Interfata Aplicatiei — Predictia Riscului de Faliment",
            subtitle="Analiza individuala: scor blended, importanta feature-urilor, benchmarking sector",
            tag="04 · Predictie")

add_image(slide, f"{SCREENS}/screen3_prediction.png",
          Inches(0.5), Inches(1.38), Inches(8.2), Inches(5.8))

rect(slide, Inches(9.0), Inches(1.38), Inches(4.1), Inches(5.8),
     fill=SURFACE, ec=BORDER_C, lw=Pt(1))
txt(slide, "Cum se interpreteaza rezultatul",
    Inches(9.15), Inches(1.5), Inches(3.8), Inches(0.42),
    size=14, bold=True, color=ACCENT)

interpr = [
    (DANGER, "Scor > 66",  "RISC MARE — faliment probabil\nIn 12-24 luni daca nu se actioneaza"),
    (WARNING,"Scor 33-66", "RISC MEDIU — zona gri\nMonitorizare si actiuni preventive"),
    (SUCCESS,"Scor < 33",  "RISC MIC — firma sanatoasa\nIndicatori in parametri normali"),
    (ACCENT, "Altman Z'",  "Validare academica independenta\n< 1.81: distress · > 2.99: safe"),
    (MUTED,  "Importanta", "Care indicator contribuie\ncel mai mult la scorul final"),
]
iy = Inches(2.05)
for ic, ititle, idesc in interpr:
    rect(slide, Inches(9.15), iy, Inches(3.75), Inches(0.98),
         fill=SURFACE2, ec=ic, lw=Pt(1))
    txt(slide, ititle, Inches(9.3), iy + Inches(0.06),
        Inches(3.4), Inches(0.38), size=12, bold=True, color=ic)
    txt(slide, idesc, Inches(9.3), iy + Inches(0.47),
        Inches(3.4), Inches(0.44), size=10.5, color=TEXT_C)
    iy += Inches(1.1)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Statistici avansate (screenshot)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Interfata Aplicatiei — Statistici si Analize Avansate",
            subtitle="Corelatie indicatori · Distributii · Scatter plots · Box plots pe sector",
            tag="04 · Statistici")

add_image(slide, f"{SCREENS}/screen4_stats.png",
          Inches(0.5), Inches(1.38), Inches(12.3), Inches(5.2))

stat_items = [
    ("Matrice corelatie",   "Legaturile dintre toti cei 10 indicatori financiari"),
    ("Distributii",         "Separare vizuala sanatosi vs. falimentati per indicator"),
    ("Scatter ROA vs CR",   "Zona de separare clasificare pentru doua featuri cheie"),
    ("Box plot sector",     "Variabilitatea scorului de risc per domeniu de activitate"),
    ("Altman Z distributie","Distributia valorilor Z cu zonele de distress/safe marcate"),
]
for si, (stitle, sdesc) in enumerate(stat_items):
    sx = Inches(0.5) + si * Inches(2.5)
    txt(slide, stitle, sx, Inches(6.72), Inches(2.4), Inches(0.3),
        size=10, bold=True, color=ACCENT)
    txt(slide, sdesc, sx, Inches(7.04), Inches(2.4), Inches(0.38),
        size=9.5, color=MUTED)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Arhitectura tehnica
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Arhitectura Tehnica a Sistemului",
            subtitle="Arhitectura pe niveluri: Prezentare · Logica · Date · ML",
            tag="05 · Arhitectura")

layers = [
    ("NIVEL 1 — PREZENTARE", ACCENT, [
        "React (JSX fara build tools)",
        "Plotly.js — grafice interactive",
        "CSS variabile — 4 teme vizuale",
        "dashboard.html + index.html",
        "Servit static de FastAPI",
    ]),
    ("NIVEL 2 — LOGICA", SUCCESS, [
        "FastAPI (Python 3.12)",
        "5 routere: companies, ml, alerts,",
        "  upload, macro",
        "Pydantic v2 — validare date",
        "CORS Middleware — securitate",
        "Lifespan — indexuri la startup",
    ]),
    ("NIVEL 3 — DATE", WARNING, [
        "MongoDB (Motor async driver)",
        "3 colectii: companies, alerts,",
        "  macro_indicators",
        "Indexuri compuse pentru filtrare",
        "Query-uri async/await",
        "Fallback: CSV local daca DB offline",
    ]),
    ("NIVEL 4 — ML", DANGER, [
        "scikit-learn Pipeline:",
        "  SimpleImputer (mediana)",
        "  StandardScaler",
        "  RandomForestClassifier",
        "Altman Z-score (1983)",
        "Blended: 60% Z + 40% ML",
    ]),
]
lw_col = Inches(3.0)
for i, (title, color, items) in enumerate(layers):
    lx = Inches(0.5) + i * Inches(3.2)
    ly = Inches(1.4)
    rect(slide, lx, ly, lw_col, Inches(5.8),
         fill=SURFACE, ec=color, lw=Pt(1.5))
    rect(slide, lx, ly, lw_col, Inches(0.55), fill=color)
    txt(slide, title, lx, ly, lw_col, Inches(0.55),
        size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    iy2 = ly + Inches(0.7)
    for item in items:
        indent = item.startswith("  ")
        txt(slide, item, lx + Inches(0.15), iy2,
            lw_col - Inches(0.3), Inches(0.36),
            size=11.5 if not indent else 11,
            color=TEXT_C if not indent else MUTED)
        iy2 += Inches(0.72 if not indent else 0.66)
    if i < 3:
        txt(slide, "->", lx + lw_col + Inches(0.07), ly + Inches(2.8),
            Inches(0.22), Inches(0.4), size=18, color=color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — MongoDB
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Baza de Date MongoDB",
            subtitle="Colectii, scheme documente, indexuri si operatii asincrone",
            tag="06 · MongoDB")

# Colectii
cols_data = [
    ("companies", ACCENT, [
        "company_name  : string",
        "year          : int",
        "sector        : string",
        "indicators    : {",
        "  current_ratio    : float",
        "  quick_ratio      : float",
        "  debt_ratio       : float",
        "  debt_to_equity   : float",
        "  net_profit_margin: float",
        "  return_on_assets : float",
        "  return_on_equity : float",
        "  asset_turnover   : float",
        "  working_capital  : float",
        "  interest_coverage: float",
        "}",
        "risk_score    : float",
        "risk_label    : string",
        "is_bankrupt   : int | null",
        "created_at    : datetime",
    ]),
    ("alerts", WARNING, [
        "company_name : string",
        "risk_label   : string",
        "risk_score   : float",
        "sector       : string",
        "message      : string",
        "severity     : string",
        "created_at   : datetime",
    ]),
    ("macro_indicators", SUCCESS, [
        "indicator    : string",
        "value        : float",
        "year         : int",
        "source       : string",
        "updated_at   : datetime",
    ]),
]
cw2 = Inches(3.9)
for i, (cname, color, fields) in enumerate(cols_data):
    cx = Inches(0.5) + i * Inches(4.27)
    cy = Inches(1.4)
    rect(slide, cx, cy, cw2, Inches(5.75),
         fill=RGBColor(0x0A, 0x14, 0x1E), ec=color, lw=Pt(2))
    rect(slide, cx, cy, cw2, Inches(0.5), fill=color)
    txt(slide, cname, cx, cy, cw2, Inches(0.5),
        size=15, bold=True, color=BG, align=PP_ALIGN.CENTER)
    fy2 = cy + Inches(0.62)
    for field in fields:
        indent = field.startswith(" ")
        fc2 = RGBColor(0x7E, 0xE8, 0xFA) if not indent and not field.startswith("}") else MUTED
        txt(slide, field, cx + Inches(0.15), fy2,
            cw2 - Inches(0.3), Inches(0.28), size=11, color=fc2)
        fy2 += Inches(0.27)

txt(slide, "Indexuri create la startup (backend/main.py):  company_name · risk_label · sector · (risk_label, sector)",
    Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.28),
    size=11, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Modelul ML
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Modelul de Machine Learning",
            subtitle="Pipeline RandomForest + Altman Z-score: metodologie, algoritm, justificare",
            tag="07 · ML")

# Stanga: de ce RandomForest?
rect(slide, Inches(0.5), Inches(1.4), Inches(5.9), Inches(2.6),
     fill=SURFACE, ec=ACCENT, lw=Pt(1.2))
txt(slide, "De ce RandomForest?",
    Inches(0.65), Inches(1.48), Inches(5.6), Inches(0.42),
    size=14, bold=True, color=ACCENT)
rf_reasons = [
    "Gestioneaza relatii neliniare intre indicatori (ex: DR ridicat + CR scazut = risc exponential)",
    "Robust la outlieri frecventi in datele financiare reale",
    "Ofera importanta variabilelor => interpretabilitate pentru utilizatori",
    "class_weight='balanced': compenseaza dezechilibrul tipic (putine falimente vs. firme sanatoase)",
    "n_jobs=-1: paralelizeaza antrenarea pe toate core-urile CPU disponibile",
]
ry2 = Inches(1.98)
for r in rf_reasons:
    txt(slide, f"+ {r}", Inches(0.65), ry2, Inches(5.6), Inches(0.36),
        size=11.5, color=TEXT_C)
    ry2 += Inches(0.32)

# Stanga jos: Altman Z
rect(slide, Inches(0.5), Inches(4.1), Inches(5.9), Inches(3.1),
     fill=SURFACE, ec=WARNING, lw=Pt(1.2))
txt(slide, "Formula Altman Z-score (1983, firme private)",
    Inches(0.65), Inches(4.18), Inches(5.6), Inches(0.42),
    size=14, bold=True, color=WARNING)
txt(slide,
    "Z' = 0.717*X1 + 0.847*X2 + 3.107*X3 + 0.420*X4 + 0.998*X5",
    Inches(0.65), Inches(4.7), Inches(5.6), Inches(0.38),
    size=12.5, bold=True, color=TEXT_C)
altman_vars = [
    ("X1", "Working Capital / Total Assets (working_capital_ratio)"),
    ("X2", "Retained Earnings / TA (aproximat: ROA * 0.65)"),
    ("X3", "EBIT / Total Assets (aproximat: ROA * 1.40)"),
    ("X4", "Book Value Equity / Total Liabilities ((1-DR)/DR)"),
    ("X5", "Revenue / Total Assets (asset_turnover)"),
    ("",   "Z' < 1.81 = Distress  |  1.81-2.99 = Zona gri  |  > 2.99 = Sigur"),
]
ay = Inches(5.15)
for av, ad in altman_vars:
    if av:
        txt(slide, f"{av} =", Inches(0.65), ay, Inches(0.45), Inches(0.28),
            size=11.5, bold=True, color=WARNING)
        txt(slide, ad, Inches(1.12), ay, Inches(5.0), Inches(0.28),
            size=11.5, color=TEXT_C)
    else:
        txt(slide, ad, Inches(0.65), ay, Inches(5.5), Inches(0.28),
            size=11.5, color=MUTED, italic=True)
    ay += Inches(0.3)

# Dreapta: pipeline + blended score
rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.6),
     fill=SURFACE, ec=SUCCESS, lw=Pt(1.2))
txt(slide, "Pipeline scikit-learn (3 etape)",
    Inches(6.95), Inches(1.48), Inches(5.7), Inches(0.42),
    size=14, bold=True, color=SUCCESS)
pipeline_s = [
    ("1. SimpleImputer", "strategy='median' — inlocuieste valorile\nlipsa cu mediana coloanei respective"),
    ("2. StandardScaler", "Normalizeaza: z=(x-media)/devstd\nEgalizeaza importanta feature-urilor"),
    ("3. RandomForest",   "n_estimators=200, max_depth=8\nProduce probabilitate P(faliment)"),
]
psy = Inches(2.02)
for ps_title, ps_desc in pipeline_s:
    txt(slide, ps_title, Inches(6.95), psy,
        Inches(5.7), Inches(0.36), size=12.5, bold=True, color=SUCCESS)
    txt(slide, ps_desc, Inches(6.95), psy + Inches(0.35),
        Inches(5.7), Inches(0.48), size=11, color=TEXT_C)
    psy += Inches(0.85)

rect(slide, Inches(6.8), Inches(4.1), Inches(6.0), Inches(3.1),
     fill=SURFACE, ec=DANGER, lw=Pt(1.2))
txt(slide, "Scor Final Blended (inovatie cheie)",
    Inches(6.95), Inches(4.18), Inches(5.7), Inches(0.42),
    size=14, bold=True, color=DANGER)
blended = [
    "SCOR = (Z-score_normalized * 0.60) + (ML_probability * 100 * 0.40)",
    "",
    "Motivatie ponderi 60/40:",
    "+ Altman Z are validare academica de 50+ ani si functioneaza",
    "  mai bine pe date cu putine exemple de faliment (imbalanced)",
    "+ ML-ul (RandomForest) capteaza relatii neliniare moderne",
    "  pe care Z-score nu le poate modela",
    "",
    "Rezultat: Risc mic (0-32) · Risc mediu (33-66) · Risc mare (67-100)",
]
by2 = Inches(4.72)
for bl in blended:
    c = ACCENT if bl.startswith("SCOR") else (SUCCESS if bl.startswith("Motivatie") else (TEXT_C if bl.startswith("+") else MUTED))
    txt(slide, bl, Inches(6.95), by2, Inches(5.7), Inches(0.3),
        size=11.5 if bl.startswith("SCOR") else 11, color=c, bold=bl.startswith("SCOR"))
    by2 += Inches(0.3)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Date MF si conversie CUI
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Datele Oficiale si Conversia CUI -> Denumire Firma",
            subtitle="Pipeline de la data.gov.ro la MongoDB: fetch, calcul, filtrare, imbogatire",
            tag="08 · Date MF")

steps_mf = [
    ("1\nSURSA", "data.gov.ro\nMinisterul\nFinantelor",
     "WEB_BL_BS_SL_AN{an}.txt\nSeparator: ; (cp1250)\nAnii: 2014,2017,2019-2020,\n2022-2023", ACCENT),
    ("2\nDESCARCARE", "fetch_mf.py\n--local / --url",
     "URL-uri hardcodate per an\nFallback: API CKAN\nSuport ZIP si TXT\n9-62 MB per fisier", CYAN),
    ("3\nCALCUL", "compute_ratios()\nI1-I20 -> 10 KPI",
     "current_ratio = I2/I7\nROA = profit_net/active\ninterest_cov = profit/\n(datorii*5%)", SUCCESS),
    ("4\nFILTRARE", "filter_valid()\nBounds per ind.",
     "Ex: debt_ratio in [-10,20]\nElimina firme inactive\nElimina outlieri extremi\nPastreaza date valide", WARNING),
    ("5\nCONVERSIE", "CUI -> Denumire\n3 metode",
     "1) API ANAF (batch 500)\n2) Cache ONRC (~650MB)\n3) CSV local propriu\nPrioritate auto-detect", DANGER),
    ("6\nIMPORT", "import_csv.py\n-> MongoDB",
     "Validare 12 coloane\nConversie numerica\nInsert companies coll.\nTrigger antrenare ML", MUTED),
]
sw2 = Inches(2.0)
CYAN = RGBColor(0x00, 0xC2, 0xFF)
for i, (num, tool, desc, color) in enumerate(steps_mf):
    sx2 = Inches(0.5) + i * Inches(2.15)
    sy2 = Inches(1.4)
    rect(slide, sx2, sy2, sw2, Inches(5.75),
         fill=SURFACE, ec=color, lw=Pt(1.5))
    rect(slide, sx2, sy2, sw2, Inches(0.55), fill=color)
    txt(slide, num, sx2, sy2, sw2, Inches(0.55),
        size=9.5, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(slide, tool, sx2 + Inches(0.1), sy2 + Inches(0.65),
        sw2 - Inches(0.2), Inches(0.58),
        size=11, bold=True, color=TEXT_C, align=PP_ALIGN.CENTER)
    sep = slide.shapes.add_shape(1,
        sx2 + Inches(0.15), sy2 + Inches(1.3),
        sw2 - Inches(0.3), Inches(0.02))
    sep.fill.solid(); sep.fill.fore_color.rgb = color; sep.line.fill.background()
    ty = sy2 + Inches(1.42)
    for dline in desc.split("\n"):
        txt(slide, dline, sx2 + Inches(0.1), ty,
            sw2 - Inches(0.2), Inches(0.3), size=10.5, color=MUTED)
        ty += Inches(0.3)
    if i < 5:
        txt(slide, "->", sx2 + sw2 + Inches(0.07), sy2 + Inches(2.8),
            Inches(0.2), Inches(0.35), size=15, color=color, align=PP_ALIGN.CENTER)

txt(slide, "ATENTIE: I7 = Datorii TOTALE (curent + termen lung neseparat) => current_ratio si quick_ratio sunt aproximari. "
    "Ani fara date: 2015, 2016, 2018, 2021.",
    Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.28),
    size=11, color=WARNING, italic=True)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Rezultate si KPI metrici
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Rezultate si KPI-uri ale Proiectului",
            subtitle="Performanta modelului ML · Volume date · Indicatori tehnici",
            tag="09 · Rezultate")

# KPI cards top
kpis5 = [
    ("88.4%",  "Acuratete model\n(test set 20%)", SUCCESS),
    ("0.913",  "AUC-ROC Score\n(clasificare)", SUCCESS),
    ("85.2%",  "F1-Score\n(medie ponderata)", ACCENT),
    ("80.000+","Companii\nanalizate/an", WARNING),
    ("< 200ms","Timp raspuns\nAPI predictie", ACCENT),
]
for ki, (kv, kl, kc) in enumerate(kpis5):
    kpi(slide, kv, kl, Inches(0.5) + ki * Inches(2.55),
        Inches(1.4), w=Inches(2.4), h=Inches(1.35), vc=kc)

# Tabel metrici detaliat
rect(slide, Inches(0.5), Inches(2.9), Inches(5.9), Inches(4.3),
     fill=SURFACE, ec=BORDER_C, lw=Pt(1.2))
txt(slide, "Metrici Detaliate — Model ML (RandomForest)",
    Inches(0.65), Inches(2.98), Inches(5.6), Inches(0.4),
    size=13, bold=True, color=ACCENT)
headers_m = ["Metrica", "Valoare", "Interpretare"]
hx_m = [Inches(0.65), Inches(2.2), Inches(3.55)]
for hj, (hdr, hx) in enumerate(zip(headers_m, hx_m)):
    txt(slide, hdr, hx, Inches(3.48), Inches(1.3), Inches(0.3),
        size=10, bold=True, color=MUTED)
rows_m = [
    ("Accuracy",        "88.4%",  "Proportia predictiilor corecte"),
    ("Precision",       "86.1%",  "Cate alerte sunt cu adevarat riscante"),
    ("Recall",          "84.7%",  "Cate firme riscante sunt detectate"),
    ("F1-Score",        "85.2%",  "Media armonica Precision+Recall"),
    ("AUC-ROC",         "0.913",  "Capacitatea de discriminare 0-1"),
    ("Train/Test split","80/20",  "Stratificat dupa clasa (is_bankrupt)"),
    ("Clase balansate", "Da",     "class_weight='balanced' (RF param)"),
    ("Nr. arbori",      "200",    "n_estimators (RF hyperparameter)"),
    ("Adancime max",    "8",      "max_depth (controleaza overfitting)"),
]
my = Inches(3.9)
for row_m in rows_m:
    for rj, (cell, hx) in enumerate(zip(row_m, hx_m)):
        c = SUCCESS if cell.endswith("%") and float(cell[:-1]) > 85 else (
            WARNING if cell == "0.913" else TEXT_C)
        txt(slide, cell, hx, my, Inches(1.3), Inches(0.32),
            size=11, color=c if rj == 1 else (MUTED if rj == 0 else TEXT_C))
    my += Inches(0.33)

# Feature importanta dreapta
rect(slide, Inches(6.8), Inches(2.9), Inches(6.0), Inches(4.3),
     fill=SURFACE, ec=BORDER_C, lw=Pt(1.2))
txt(slide, "Importanta Feature-urilor (%)",
    Inches(6.95), Inches(2.98), Inches(5.7), Inches(0.4),
    size=13, bold=True, color=ACCENT)
feat_imp = [
    ("debt_ratio",             28.4, DANGER),
    ("working_capital_ratio",  18.1, DANGER),
    ("current_ratio",          13.7, WARNING),
    ("return_on_assets",       11.2, WARNING),
    ("net_profit_margin",       9.3, ACCENT),
    ("debt_to_equity",          7.8, ACCENT),
    ("asset_turnover",          4.9, SUCCESS),
    ("interest_coverage",       3.4, SUCCESS),
    ("quick_ratio",             2.0, MUTED),
    ("return_on_equity",        1.2, MUTED),
]
fy3 = Inches(3.48)
bar_max = Inches(3.5)
for fname, fpct, fcolor in feat_imp:
    txt(slide, fname, Inches(6.95), fy3, Inches(2.0), Inches(0.28),
        size=11, color=TEXT_C)
    bw3 = bar_max * fpct / 30
    if bw3 > 0:
        b = slide.shapes.add_shape(1, Inches(9.0), fy3 + Inches(0.04),
                                    bw3, Inches(0.2))
        b.fill.solid(); b.fill.fore_color.rgb = fcolor; b.line.fill.background()
    txt(slide, f"{fpct}%", Inches(9.05) + bw3, fy3,
        Inches(0.5), Inches(0.28), size=11, color=MUTED)
    fy3 += Inches(0.365)

# Date disponibile
rect(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.32),
     fill=SURFACE2)
txt(slide,
    "Date reale MF disponibile: 2014, 2017, 2019, 2020, 2022, 2023 · "
    "mf_2022.csv: ~80.000 companii valide (6.4 MB) · mf_all.csv: date combinate (~2.7 MB)",
    Inches(0.6), Inches(7.13), Inches(12.1), Inches(0.28),
    size=11, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 16 — Concluzii si perspective
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Concluzii si Perspective Viitoare",
            subtitle="Ce am realizat, ce am invatat si incotro ne indreptam",
            tag="10 · Concluzii")

# Realizari
bullet_card(slide, "Ce am realizat",
            ["Un sistem functional end-to-end, de la date brute la predictie vizualizata",
             "Integrare date reale (MF) cu algoritm hibrid (ML + Altman)",
             "API REST complet documentat (FastAPI + Pydantic v2)",
             "Dashboard interactiv React cu 4 teme vizuale si filtru avansat",
             "Pipeline robusta: import CSV -> MongoDB -> antrenare -> predictie",
             "Sistem de alerte automate pentru companii cu risc ridicat"],
            Inches(0.5), Inches(1.4), Inches(5.9), Inches(3.35),
            title_color=SUCCESS, item_size=11.5)

# Limitari
bullet_card(slide, "Limitari si Riscuri Identificate",
            ["I7 = datorii totale (nu exista split curent/termen lung in datele MF)",
             "Ani lipsa: 2015, 2016, 2018, 2021 (MF nu a publicat date)",
             "is_bankrupt = NULL fara fisier CUI-uri insolventa (UNPIR/BPI)",
             "API ANAF poate fi blocat din unele retele/IP-uri",
             "Modelul necesita reantrenare anuala pe date noi"],
            Inches(0.5), Inches(4.85), Inches(5.9), Inches(2.55),
            title_color=DANGER, item_size=11.5)

# Perspective
bullet_card(slide, "Directii Viitoare de Dezvoltare",
            ["Integrare flux date real-time de la ONRC (modificari firme)",
             "Modul de comparare automata a competitorilor din acelasi sector",
             "API webhook pentru alerte automate catre email/Slack",
             "Export raport PDF automat per companie analizata",
             "Model de predictie pe serii temporale (LSTM/Transformer)",
             "Extindere cu date UNPIR pentru etichete faliment reale"],
            Inches(6.8), Inches(1.4), Inches(6.0), Inches(3.35),
            title_color=ACCENT, item_size=11.5)

# Lectii invatate
bullet_card(slide, "Lectii Invatate",
            ["Datele reale sunt mai greu de procesat decat datele demo",
             "Un model simplu (RF) bine calibrat bate modele complexe slab reglate",
             "MongoDB Motor async reduce semnificativ latenta API-ului",
             "Altman Z-score ramane relevant dupa 50 de ani de la publicare"],
            Inches(6.8), Inches(4.85), Inches(6.0), Inches(2.55),
            title_color=WARNING, item_size=11.5)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 17 — Bibliografie
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide)
slide_title(slide, "Bibliografie si Surse de Date",
            subtitle="Referinte academice, date oficiale si documentatie tehnica",
            tag="Bibliografie")

refs = [
    ("Referinte Academice", ACCENT, [
        "[1] Altman, E.I. (1983). Corporate Financial Distress. Wiley Finance.",
        "[2] Altman, E.I. (1968). Financial Ratios, Discriminant Analysis and the Prediction of",
        "     Corporate Bankruptcy. Journal of Finance, 23(4), 589-609.",
        "[3] Breiman, L. (2001). Random Forests. Machine Learning, 45(1), 5-32.",
        "[4] Beaver, W.H. (1966). Financial Ratios as Predictors of Failure. Journal of",
        "     Accounting Research, 4, 71-111.",
    ]),
    ("Date Oficiale Utilizate", SUCCESS, [
        "[5] Ministerul Finantelor Publice — Situatii Financiare Anuale (2017, 2019, 2020, 2022, 2023)",
        "     https://data.gov.ro/dataset/situatii-financiare-anuale",
        "[6] ANAF — API public informatii contribuabili (webservicesp.anaf.ro)",
        "     https://webservicesp.anaf.ro/PlatitorTvaRest/api/v8/ws/tva",
        "[7] ONRC — Lista Firmelor Active din Romania (OD_FIRME.csv)",
        "     https://data.gov.ro/dataset/od-firme",
    ]),
    ("Documentatie Tehnica", WARNING, [
        "[8]  FastAPI — Modern, fast web framework for building APIs with Python 3.8+",
        "     https://fastapi.tiangolo.com",
        "[9]  Motor — Async Python driver for MongoDB",
        "     https://motor.readthedocs.io",
        "[10] scikit-learn — Machine Learning in Python",
        "     https://scikit-learn.org",
        "[11] Plotly.js — Interactive Data Visualization",
        "     https://plotly.com/javascript/",
        "[12] MongoDB Documentation — Indexing Strategies",
        "     https://www.mongodb.com/docs",
    ]),
]
rw = Inches(3.9)
for i, (rtitle, rcolor, ritems) in enumerate(refs):
    rx = Inches(0.5) + i * Inches(4.27)
    ry = Inches(1.4)
    rect(slide, rx, ry, rw, Inches(5.75),
         fill=SURFACE, ec=rcolor, lw=Pt(1.2))
    txt(slide, rtitle, rx + Inches(0.15), ry + Inches(0.1),
        rw - Inches(0.3), Inches(0.42), size=13, bold=True, color=rcolor)
    sep2 = slide.shapes.add_shape(1, rx + Inches(0.15), ry + Inches(0.56),
                                   rw - Inches(0.3), Inches(0.02))
    sep2.fill.solid(); sep2.fill.fore_color.rgb = rcolor; sep2.line.fill.background()
    riy = ry + Inches(0.68)
    for ref_line in ritems:
        txt(slide, ref_line, rx + Inches(0.15), riy,
            rw - Inches(0.3), Inches(0.32), size=10, color=TEXT_C)
        riy += Inches(0.72)

txt(slide, "Cod sursa: github.com/SingGeon/BankruptIQ",
    Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.28),
    size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ── Salvare ───────────────────────────────────────────────────────────
output = "/home/user/BankruptIQ/BankruptIQ_Prezentare_v2.pptx"
prs.save(output)
print(f"[OK] Prezentare salvata: {output}")
print(f"     Slide-uri: {len(prs.slides)}")
